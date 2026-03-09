import io
import os
import logging
from fastapi import HTTPException, UploadFile
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Pt

logger = logging.getLogger(__name__)

async def convert_pdf_to_pptx(file: UploadFile):
    """
    Convert a PDF file to a PowerPoint (.pptx) presentation.
    - Each PDF page becomes a high-resolution slide.
    - Slide dimensions match PDF page dimensions for perfect layout.
    """
    if not file.filename or not file.filename.lower().endswith(".pdf"):
        raise HTTPException(
            status_code=400,
            detail="Only PDF files are accepted. Please upload a .pdf file.",
        )

    logger.info(f"Converting to PowerPoint: {file.filename}")

    pdf_bytes = await file.read()
    if len(pdf_bytes) == 0:
        raise HTTPException(status_code=400, detail="Uploaded file is empty.")

    try:
        # Open PDF from bytes
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        prs = Presentation()
        
        # Remove the first default slide created by 'Presentation()'
        # (Though prs.slides is initially empty if we don't pick a template)
        
        for i in range(len(doc)):
            page = doc.load_page(i)
            # Match slide size to PDF page (in Points)
            # PowerPoint width/height are in 'English Metric Units' (EMUs)
            # 1 Point = 12,700 EMUs
            width_pts = page.rect.width
            height_pts = page.rect.height
            
            prs.slide_width = Pt(width_pts)
            prs.slide_height = Pt(height_pts)
            
            # Add a blank slide layout
            blank_slide_layout = prs.slide_layouts[6] 
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Render page to a high-res image (300 DPI)
            pix = page.get_pixmap(matrix=fitz.Matrix(300/72, 300/72))
            img_data = pix.tobytes("png")
            
            # Insert image as a full-slide-matching picture
            img_stream = io.BytesIO(img_data)
            slide.shapes.add_picture(img_stream, 0, 0, width=Pt(width_pts), height=Pt(height_pts))

        # Serialize to bytes
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_bytes = pptx_buffer.getvalue()
        doc.close()

    except Exception as e:
        logger.error(f"PowerPoint conversion failed: {e}")
        raise HTTPException(status_code=500, detail=f"Conversion failed: {str(e)}")

    base_name = os.path.splitext(file.filename)[0]
    output_filename = f"{base_name}.pptx"
    
    return pptx_bytes, output_filename
